from io import BytesIO
from docx import Document
from docx.shared import Pt
from pptx import Presentation
from pptx.util import Pt, Inches


# ---------------- DOCX (Beautiful & Clean) ---------------- #
def assemble_docx(project, sections):
    doc = Document()

    # Title
    title = doc.add_heading(project.title or "Document", level=1)
    for run in title.runs:
        run.font.name = "Segoe UI"
        run.font.size = Pt(28)

    if project.topic:
        topic_para = doc.add_paragraph(f"Topic: {project.topic}")
        topic_para.runs[0].font.size = Pt(13)
        topic_para.runs[0].font.name = "Segoe UI"
        doc.add_paragraph("")

    # Sections
    for s in sorted(sections, key=lambda x: x.order):
        heading = doc.add_heading(s.title or "Section", level=2)
        for run in heading.runs:
            run.font.name = "Segoe UI"
            run.font.size = Pt(20)

        lines = (s.content or "").split("\n")
        for line in lines:
            line = line.strip()
            if not line:
                continue

            # Bullet formatting if PPT bullet-style content comes from LLM
            if line.startswith("-") or line.startswith("•"):
                p = doc.add_paragraph(line.strip("-• "), style="List Bullet")
                p.runs[0].font.name = "Segoe UI"
                p.runs[0].font.size = Pt(11)
            else:
                p = doc.add_paragraph(line)
                p.runs[0].font.name = "Segoe UI"
                p.runs[0].font.size = Pt(12)

        doc.add_paragraph("")

    bio = BytesIO()
    doc.save(bio)
    bio.seek(0)
    return bio



# ---------------- PPTX (Professional + Business Style) ---------------- #
def assemble_pptx(project, sections):
    prs = Presentation()

    # Title Slide
    title_slide = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide)
    slide.shapes.title.text = project.title or "Untitled Presentation"

    if project.topic and len(slide.placeholders) > 1:
        slide.placeholders[1].text = f"Topic: {project.topic}"

    # Style for bullets
    font_size = Pt(24)
    font_name = "Segoe UI"

    # Content Slides
    for s in sorted(sections, key=lambda x: x.order):
        layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = s.title or "Slide"

        textbox = slide.shapes.placeholders[1].text_frame
        textbox.clear()

        lines = [l.strip() for l in (s.content or "").split("\n") if l.strip()]
        for i, line in enumerate(lines):
            if i == 0:
                p = textbox.paragraphs[0]
            else:
                p = textbox.add_paragraph()
            p.text = line
            p.level = 0
            p.font.size = font_size
            p.font.name = font_name

    # Export to bytes
    bio = BytesIO()
    prs.save(bio)
    bio.seek(0)
    return bio
